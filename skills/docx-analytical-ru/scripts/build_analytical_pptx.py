#!/usr/bin/env python3
"""
Шаблон презентации для аналитических комплектов (docx-analytical-ru).

Содержит стили и хелперы. Для конкретного проекта — скопировать в свой git-репозиторий
и задать контент слайдов (переименовать в build_<project>_pptx.py и наполнить своими данными).

Usage:
  python3 build_analytical_pptx.py OUTPUT.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x38, 0x64)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x36, 0x45, 0x4F)
FONT = "Times New Roman"


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title(slide, text: str, *, dark: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE if dark else NAVY


def add_body(slide, lines: list[str], top: float = 1.5, size: int = 18, *, light: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    color = WHITE if light else CHARCOAL
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_card(slide, left, top, width, height, title, body, fill=NAVY):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(16)
    p1.font.name = FONT
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.font.name = FONT
    p2.font.color.rgb = ICE
    p2.alignment = PP_ALIGN.CENTER


def build_minimal_demo(output: Path) -> None:
    """Минимальный пример: титул + резюме + следующие шаги."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Заголовок презентации"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title(s, "Резюме")
    add_body(s, ["• Замените этот шаблон контентом проекта", "• Только русский текст в слайдах"], top=1.7)

    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    add_title(s, "Следующие шаги", dark=True)
    add_body(s, ["1. …", "2. …"], top=2.2, size=22, light=True)

    prs.save(str(output))
    print(f"OK {output.name} (demo template — customize for project)")


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("Презентация_аналитика_шаблон.pptx")
    build_minimal_demo(out)
